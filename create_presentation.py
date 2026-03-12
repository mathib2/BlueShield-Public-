"""
Generate the BlueShield NABC Presentation — Polished version with logo and assets.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import os

PROJ = Path(__file__).parent
ASSETS = PROJ / "assets"
LOGO = str(ASSETS / "logo.png")
APP_MOCKUP = str(ASSETS / "app_mockup.png")
BT_HACKING = str(ASSETS / "bt_hacking.png")
OUTPUT = str(PROJ / "BlueShield_NABC.pptx")

# ── Color Palette ────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0B, 0x0F, 0x19)
BG_CARD      = RGBColor(0x12, 0x1A, 0x2E)
BG_ACCENT    = RGBColor(0x18, 0x22, 0x3A)
BLUE         = RGBColor(0x1A, 0x6B, 0xB5)   # brand blue from logo
CYAN         = RGBColor(0x00, 0xD4, 0xFF)
GREEN        = RGBColor(0x00, 0xE6, 0x76)
RED          = RGBColor(0xFF, 0x3B, 0x5C)
ORANGE       = RGBColor(0xFF, 0x9F, 0x43)
PURPLE       = RGBColor(0x8B, 0x5C, 0xF6)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LT      = RGBColor(0xA0, 0xAE, 0xC0)
GRAY_DIM     = RGBColor(0x5A, 0x6A, 0x80)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5


def bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def txt(slide, x, y, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tb


def bullets(slide, x, y, w, h, items, size=14, color=GRAY_LT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        from pptx.oxml.ns import qn
        pPr = p._p.get_or_add_pPr()
        for ch in list(pPr):
            tag = ch.tag.split('}')[-1] if '}' in ch.tag else ch.tag
            if tag.startswith('bu'):
                pPr.remove(ch)
        pPr.append(pPr.makeelement(qn('a:buFont'), {'typeface': 'Arial'}))
        pPr.append(pPr.makeelement(qn('a:buChar'), {'char': '\u25CF'}))
    return tb


def stat(slide, x, y, w, h, number, label, color=CYAN):
    rect(slide, x, y, w, h, BG_CARD)
    txt(slide, x, y + 0.15, w, 0.8, number, size=42, color=color, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x, y + 1.0, w, 0.4, label, size=12, color=GRAY_LT, align=PP_ALIGN.CENTER)


def card(slide, x, y, w, h, title, body_lines, accent=CYAN):
    rect(slide, x, y, w, h, BG_CARD)
    rect(slide, x, y, w, 0.06, accent)
    txt(slide, x + 0.3, y + 0.2, w - 0.6, 0.35, title, size=16, color=WHITE, bold=True)
    if body_lines:
        txt(slide, x + 0.3, y + 0.7, w - 0.6, h - 0.9, "\n".join(body_lines), size=12, color=GRAY_LT)


def footer(slide):
    rect(slide, 0, SH - 0.06, SW, 0.06, BLUE)


def logo_small(slide, x, y, h=0.6):
    """Add logo image."""
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(x), Inches(y), height=Inches(h))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
rect(s, 0, 3.15, SW, 0.012, CYAN)

# Logo centered
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(5.4), Inches(0.4), height=Inches(2.2))

txt(s, 1.0, 2.6, 11.3, 0.5, "BLUETOOTH SECURITY MONITORING PLATFORM", size=14, color=CYAN, align=PP_ALIGN.CENTER)
txt(s, 1.0, 3.6, 11.3, 0.7, "Protecting Your Wireless Environment", size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, 1.0, 4.4, 11.3, 0.5, "Research-Grade Sniffing & Jamming Device  |  Raspberry Pi Powered", size=14, color=GRAY_LT, align=PP_ALIGN.CENTER)

# Team
txt(s, 1.0, 5.5, 11.3, 0.35, "Mathias Vera   \u2022   Daniel Halbleib Jr   \u2022   Andrew Sauls", size=15, color=GRAY_DIM, align=PP_ALIGN.CENTER)
txt(s, 1.0, 5.9, 11.3, 0.3, "NABC Presentation  |  2026", size=11, color=GRAY_DIM, align=PP_ALIGN.CENTER)
footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
logo_small(s, 0.5, 0.3)

txt(s, 1.3, 0.35, 5, 0.5, "THE PROBLEM", size=34, color=WHITE, bold=True)
rect(s, 1.3, 0.95, 1.5, 0.04, CYAN)

txt(s, 0.8, 1.3, 7.5, 1.0,
    "Bluetooth devices are everywhere in modern workplaces \u2014 headphones, keyboards, "
    "speakers, fitness trackers. Yet most organizations have ZERO visibility into what "
    "Bluetooth devices are operating in their space.",
    size=15, color=GRAY_LT)

# Hacking illustration
if os.path.exists(BT_HACKING):
    s.shapes.add_picture(BT_HACKING, Inches(8.5), Inches(0.8), width=Inches(4.2))

threats = [
    ("Eavesdropping", "Attackers intercept audio from\nBT headsets and speakers used\nin sensitive meetings", RED),
    ("Unauthorized Access", "Rogue devices pair with company\nhardware, exfiltrating data\nthrough BT connections", ORANGE),
    ("BlueBorne Attacks", "Unpatched BT stacks are vulnerable\nto remote code execution without\nany user interaction", PURPLE),
]
for i, (title, body, color) in enumerate(threats):
    cx = 0.8 + i * 4.0
    rect(s, cx, 3.4, 3.6, 2.8, BG_CARD)
    rect(s, cx, 3.4, 3.6, 0.06, color)
    txt(s, cx + 0.3, 3.7, 3.0, 0.35, title, size=17, color=color, bold=True)
    txt(s, cx + 0.3, 4.2, 3.0, 1.6, body, size=13, color=GRAY_LT)
footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — TARGET MARKET
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
logo_small(s, 0.5, 0.3)

txt(s, 1.3, 0.35, 5, 0.5, "WHO'S AT RISK?", size=34, color=WHITE, bold=True)
rect(s, 1.3, 0.95, 1.5, 0.04, CYAN)
txt(s, 0.8, 1.2, 8, 0.4, "Any organization where Bluetooth devices operate near sensitive data or conversations.", size=14, color=GRAY_LT)

markets = [
    ("Corporate Offices", ["Wireless headphones & speakers", "BT keyboards & mice", "Conference room devices", "BYOD personal phones"], CYAN),
    ("Finance & Legal", ["Client-confidential calls", "Sensitive document transfers", "Trading floor comms", "Regulatory compliance"], GREEN),
    ("Healthcare", ["HIPAA-protected conversations", "Medical IoT devices (BLE)", "Patient data on wireless", "Telemedicine endpoints"], ORANGE),
    ("Government & Defense", ["Classified facility monitoring", "SCIF perimeter security", "Counter-intelligence ops", "Insider threat detection"], RED),
]
for i, (title, items, color) in enumerate(markets):
    cx = 0.5 + i * 3.15
    rect(s, cx, 2.0, 2.95, 4.8, BG_CARD)
    rect(s, cx, 2.0, 2.95, 0.06, color)
    txt(s, cx + 0.2, 2.2, 2.55, 0.35, title, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    rect(s, cx + 0.3, 2.65, 2.35, 0.015, GRAY_DIM)
    bullets(s, cx + 0.2, 2.85, 2.55, 3.5, items, size=12, color=GRAY_LT)
footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — OUR APPROACH
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
logo_small(s, 0.5, 0.3)

txt(s, 1.3, 0.35, 5, 0.5, "OUR APPROACH", size=34, color=WHITE, bold=True)
rect(s, 1.3, 0.95, 1.5, 0.04, CYAN)

steps = [
    ("1", "DETECT", "Dedicated Raspberry Pi monitor\nplaced in key areas. Passive\nBLE + Classic BT scanning\nvia HCI interface.", CYAN),
    ("2", "ANALYZE", "Real-time pattern detection.\nIdentify unknown devices,\nanomalous behavior, and\npotential attack signatures.", PURPLE),
    ("3", "RESPOND", "Instant admin alerts via\nweb dashboard. Optional\nreactive jamming. Full audit\nlogging for incident response.", RED),
]
for i, (num, title, body, color) in enumerate(steps):
    cx = 0.8 + i * 4.1
    rect(s, cx, 1.4, 3.7, 3.2, BG_CARD)
    rect(s, cx, 1.4, 0.08, 3.2, color)
    txt(s, cx + 0.3, 1.55, 0.8, 0.6, num, size=36, color=color, bold=True)
    txt(s, cx + 0.9, 1.6, 2.5, 0.45, title, size=22, color=WHITE, bold=True)
    txt(s, cx + 0.3, 2.3, 3.1, 2.0, body, size=14, color=GRAY_LT)
    if i < 2:
        txt(s, cx + 3.7, 2.6, 0.4, 0.5, "\u25B6", size=24, color=CYAN, align=PP_ALIGN.CENTER)

# Tech stack
rect(s, 0.8, 5.0, 11.7, 1.8, BG_ACCENT)
txt(s, 1.2, 5.1, 3, 0.35, "TECH STACK", size=14, color=CYAN, bold=True)
techs = ["Raspberry Pi 4/5", "BlueZ (Linux BT stack)", "Python + Bleak (BLE)",
         "HCItool + HCIdump", "Flask Web Dashboard", "JSON Audit Logging"]
for i, item in enumerate(techs):
    col, row = i % 3, i // 3
    txt(s, 1.2 + col * 4.0, 5.55 + row * 0.4, 3.5, 0.35, f"\u25A0  {item}", size=12, color=GRAY_LT)
footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DASHBOARD PREVIEW
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
logo_small(s, 0.5, 0.3)

txt(s, 1.3, 0.35, 5, 0.5, "WEB DASHBOARD", size=34, color=WHITE, bold=True)
rect(s, 1.3, 0.95, 1.5, 0.04, CYAN)

# Terminal simulation box
tx, ty, tw, th = 0.8, 1.3, 8.0, 5.5
rect(s, tx, ty, tw, 0.4, RGBColor(0x2D, 0x33, 0x3F))
rect(s, tx, ty + 0.4, tw, th - 0.4, RGBColor(0x0A, 0x0E, 0x14))

for j, c in enumerate([RED, ORANGE, GREEN]):
    btn = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(tx + 0.2 + j * 0.3), Inches(ty + 0.1), Inches(0.18), Inches(0.18))
    btn.fill.solid()
    btn.fill.fore_color.rgb = c
    btn.line.fill.background()

txt(s, tx + 1.5, ty + 0.05, 5, 0.35, "BlueShield Web Dashboard  |  http://localhost:8080", size=10, color=GRAY_DIM, align=PP_ALIGN.CENTER)

lines = [
    ("  BLUESHIELD v0.1.0  |  Bluetooth Security Monitor  |  2026-03-11", CYAN),
    ("  Scanner: SCANNING    Jammer: OFF    Scans: 47    Interface: hci0", GRAY_LT),
    ("", WHITE),
    ("  \u250C\u2500 DEVICE SUMMARY \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510", GRAY_DIM),
    ("  \u2502 Total: 13    Known: 5    Unknown: 8    Alerts: 3    \u2502", WHITE),
    ("  \u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518", GRAY_DIM),
    ("", WHITE),
    ("  ADDRESS            NAME             TYPE    RSSI     ALERT", CYAN),
    ("  \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500", GRAY_DIM),
    ("  C9:FA:86:C0:DC:4E  Amazfit Band 5   ble    -83dBm   WARNING", ORANGE),
    ("  BE:67:01:F3:0C:21  ELK-BLEDOM01     ble    -62dBm   WARNING", ORANGE),
    ("  EC:47:0C:4B:D1:3F  Platco PR1       ble    -65dBm   CRITICAL", RED),
    ("  63:BD:DD:71:5F:63  Apple Device     ble    -49dBm   OK", GREEN),
    ("  EC:B8:3B:E7:7A:E7  Apple Device     ble    -47dBm   OK", GREEN),
]
for i, (text, color) in enumerate(lines):
    txt(s, tx + 0.15, ty + 0.5 + i * 0.3, tw - 0.3, 0.3, text, size=10, color=color, font="Consolas")

# App mockup on the right
if os.path.exists(APP_MOCKUP):
    s.shapes.add_picture(APP_MOCKUP, Inches(9.2), Inches(1.3), height=Inches(5.2))

footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — REVENUE MODEL
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
logo_small(s, 0.5, 0.3)

txt(s, 1.3, 0.35, 5, 0.5, "REVENUE MODEL", size=34, color=WHITE, bold=True)
rect(s, 1.3, 0.95, 1.5, 0.04, CYAN)

tiers = [
    ("Research License", "$0", "Open-source for academic\nand security research.\nFull scanner + dashboard.\nCommunity supported.", CYAN),
    ("Professional", "$299/unit", "Pre-configured RPi hardware\n+ software bundle.\nTechnical support.\n1 year of updates.", GREEN),
    ("Enterprise", "$999/yr", "Multi-device deployment.\nCentral management.\nCustom alert integrations.\nDedicated support.", PURPLE),
    ("Managed Service", "Custom", "Fully managed monitoring.\n24/7 SOC integration.\nIncident response.\nCompliance reporting.", ORANGE),
]
for i, (name, price, desc, color) in enumerate(tiers):
    cx = 0.5 + i * 3.15
    rect(s, cx, 1.5, 2.95, 4.5, BG_CARD)
    rect(s, cx, 1.5, 2.95, 0.06, color)
    txt(s, cx + 0.2, 1.75, 2.55, 0.3, name, size=13, color=GRAY_LT, bold=True, align=PP_ALIGN.CENTER)
    txt(s, cx + 0.2, 2.15, 2.55, 0.6, price, size=34, color=color, bold=True, align=PP_ALIGN.CENTER)
    rect(s, cx + 0.3, 2.9, 2.35, 0.015, GRAY_DIM)
    txt(s, cx + 0.2, 3.1, 2.55, 2.5, desc, size=12, color=GRAY_LT, align=PP_ALIGN.CENTER)
footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — BENEFITS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
logo_small(s, 0.5, 0.3)

txt(s, 1.3, 0.35, 5, 0.5, "THE BENEFITS", size=34, color=WHITE, bold=True)
rect(s, 1.3, 0.95, 1.5, 0.04, CYAN)

stats_data = [
    ("<5s", "Alert Response", CYAN),
    ("24/7", "Monitoring", GREEN),
    ("100%", "BLE Detection", PURPLE),
    ("<$50", "Hardware Cost", ORANGE),
]
for i, (n, lbl, c) in enumerate(stats_data):
    stat(s, 0.5 + i * 3.15, 1.4, 2.95, 1.4, n, lbl, c)

benefits = [
    ("Real-Time Monitoring", "Continuous passive scanning detects every Bluetooth\ndevice the moment it appears. No gaps, no blind spots."),
    ("Privacy Protection", "Identify devices that could eavesdrop on sensitive\nconversations or exfiltrate data wirelessly."),
    ("Audit & Compliance", "Complete JSON logging of every device, alert, and\naction. Ready for auditors and incident response."),
    ("Easy Deployment", "Plug-and-play RPi hardware. Web dashboard runs in\nany browser. No complex setup or training needed."),
]
for i, (title, body) in enumerate(benefits):
    col, row = i % 2, i // 2
    bx = 0.5 + col * 6.3
    by = 3.2 + row * 2.0
    rect(s, bx, by, 6.0, 1.7, BG_CARD)
    rect(s, bx, by, 0.08, 1.7, CYAN)
    txt(s, bx + 0.3, by + 0.15, 5.4, 0.35, title, size=16, color=WHITE, bold=True)
    txt(s, bx + 0.3, by + 0.55, 5.4, 1.0, body, size=12, color=GRAY_LT)
footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CONSTRAINTS & COMPETITION
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
logo_small(s, 0.5, 0.3)

txt(s, 1.3, 0.35, 5, 0.5, "CONSTRAINTS & COMPETITION", size=34, color=WHITE, bold=True)
rect(s, 1.3, 0.95, 1.5, 0.04, CYAN)

txt(s, 0.8, 1.3, 5.5, 0.35, "CONSTRAINTS", size=17, color=ORANGE, bold=True)
constraints = [
    "FCC regulations restrict jamming in production use",
    "BT 5.x frequency hopping limits passive sniffing",
    "RPi range limited to ~30m without external antenna",
    "Requires root access for HCI-level operations",
    "CPU limits on high-density BLE environments",
]
bullets(s, 0.8, 1.8, 5.5, 2.8, constraints, size=13, color=GRAY_LT)

txt(s, 7.0, 1.3, 5.5, 0.35, "COMPETITIVE LANDSCAPE", size=17, color=PURPLE, bold=True)
competitors = [
    ("Kismet", "General wireless monitor. No BT-specific\nalerting. Complex setup required."),
    ("Ubertooth One", "Powerful BT hardware. $120+ per unit.\nRequires deep technical expertise."),
    ("Enterprise WIDS", "Cisco/Aruba solutions. $10K+ per\ndeployment. Not Bluetooth-focused."),
]
for i, (name, desc) in enumerate(competitors):
    cy = 1.9 + i * 1.55
    rect(s, 7.0, cy, 5.5, 1.3, BG_CARD)
    rect(s, 7.0, cy, 0.08, 1.3, PURPLE)
    txt(s, 7.3, cy + 0.1, 4.9, 0.3, name, size=14, color=WHITE, bold=True)
    txt(s, 7.3, cy + 0.45, 4.9, 0.7, desc, size=12, color=GRAY_LT)

# Advantage bar
rect(s, 0.8, 5.6, 11.7, 1.2, BG_ACCENT)
rect(s, 0.8, 5.6, 11.7, 0.05, GREEN)
txt(s, 1.2, 5.75, 3, 0.35, "BLUESHIELD ADVANTAGE", size=13, color=GREEN, bold=True)
txt(s, 1.2, 6.15, 10.5, 0.4,
    "Low-cost (<$50 hardware)  \u25CF  Purpose-built for BT security  \u25CF  "
    "Research-grade + commercial path  \u25CF  Open architecture", size=12, color=GRAY_LT)
footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — WHAT WE NEED
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
logo_small(s, 0.5, 0.3)

txt(s, 1.3, 0.35, 5, 0.5, "WHAT WE NEED", size=34, color=WHITE, bold=True)
rect(s, 1.3, 0.95, 1.5, 0.04, CYAN)

needs = [
    ("Hardware", ["Raspberry Pi 4/5 units", "External BT antennas", "USB BT 5.0+ adapters", "Protective deployment cases"], CYAN, "$500"),
    ("Software & Tools", ["Cloud hosting for dashboard", "Security testing lab access", "Code signing certificates", "CI/CD pipeline"], PURPLE, "$300"),
    ("Research & Testing", ["Controlled RF test environment", "Various BT devices for testing", "Security audit / pentest review", "Academic publication support"], GREEN, "$400"),
    ("Go-to-Market", ["Product packaging & branding", "Demo units for conferences", "Marketing website", "Legal review (FCC compliance)"], ORANGE, "$800"),
]
for i, (title, items, color, cost) in enumerate(needs):
    col, row = i % 2, i // 2
    nx = 0.5 + col * 6.3
    ny = 1.4 + row * 2.8
    rect(s, nx, ny, 6.0, 2.5, BG_CARD)
    rect(s, nx, ny, 0.08, 2.5, color)
    txt(s, nx + 0.3, ny + 0.12, 4.0, 0.35, title, size=17, color=WHITE, bold=True)
    txt(s, nx + 4.5, ny + 0.12, 1.3, 0.35, cost, size=17, color=color, bold=True, align=PP_ALIGN.RIGHT)
    bullets(s, nx + 0.3, ny + 0.6, 5.4, 1.8, items, size=12, color=GRAY_LT)

rect(s, 4.5, 7.0 - 0.55, 4.3, 0.5, BG_ACCENT)
txt(s, 4.5, 7.0 - 0.5, 4.3, 0.4, "TOTAL ASK:  ~$2,000", size=18, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CLOSING
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
rect(s, 0, 3.0, SW, 0.012, CYAN)

if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(5.4), Inches(0.3), height=Inches(2.2))

txt(s, 1.0, 2.6, 11.3, 0.5, "SECURING THE INVISIBLE WIRELESS LAYER", size=14, color=CYAN, align=PP_ALIGN.CENTER)
txt(s, 1.0, 3.5, 11.3, 0.6, "Thank You", size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, 1.0, 4.3, 11.3, 0.4, "Questions?", size=20, color=GRAY_LT, align=PP_ALIGN.CENTER)
txt(s, 1.0, 5.3, 11.3, 0.35, "Mathias Vera   \u2022   Daniel Halbleib Jr   \u2022   Andrew Sauls", size=15, color=GRAY_DIM, align=PP_ALIGN.CENTER)
txt(s, 1.0, 5.8, 11.3, 0.3, "github.com/your-repo  |  blueshield@security", size=11, color=GRAY_DIM, align=PP_ALIGN.CENTER)
footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"Presentation saved: {OUTPUT}")
